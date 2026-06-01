"""PPTX parser using python-pptx.

Migrated from LightRAG native PPTX parser.  Extracts text from all
shapes in all slides.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from aurora_ext.rag.parser.base import BaseParser, ParseResult

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """Parser for .pptx files using ``python-pptx``."""

    _EXTENSIONS = {"pptx"}

    @property
    def supported_extensions(self) -> set[str]:
        return self._EXTENSIONS

    async def parse(self, file_path: str | Path, **kwargs: Any) -> ParseResult:
        from pptx import Presentation

        path = Path(file_path)

        try:
            prs = Presentation(str(path))
        except Exception as exc:
            raise ValueError(f"Cannot open PPTX: {path}") from exc

        slide_texts: list[str] = []
        for slide in prs.slides:
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slide_texts.append("\n".join(texts))

        full_text = "\n\n".join(slide_texts)

        return ParseResult(
            text=full_text,
            file_path=str(path),
            file_type="pptx",
            metadata={
                "slide_count": len(prs.slides),
                "char_count": len(full_text),
            },
        )
